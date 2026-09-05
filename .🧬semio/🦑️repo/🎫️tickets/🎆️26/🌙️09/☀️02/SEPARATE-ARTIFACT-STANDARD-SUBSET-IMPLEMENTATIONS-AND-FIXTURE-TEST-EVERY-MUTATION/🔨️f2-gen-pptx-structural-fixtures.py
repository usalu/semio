#!/usr/bin/env python3
"""🔨️ F2 — generates before/after PPTX fixtures for s.stdio.pptx@ecma-376/strict's 10 and
@ecma-376/transitional's 6 unfixtured mutations. Same pattern as the docx structural generator
(🔨️f2-gen-docx-structural-fixtures.py) and built on the same shared helpers
(🔨️f2_ooxml_common.py): set-snapshot is genuinely third-party-generated (both sides real
python-pptx 1.0.2 output); every OOXML-strict-vs-transitional structural mutation is handcrafted on
top of a genuine python-pptx base package, since python-pptx's public API exposes no `conformance`
attribute, namespace-URI, or VML/AlternateContent editing surface.
Idempotent: safe to re-run.
"""
import io
import json
import sys
from pathlib import Path

from pptx import Presentation

import importlib.util

_spec = importlib.util.spec_from_file_location("f2_ooxml_common", Path(__file__).parent / "🔨️f2_ooxml_common.py")
_mod = importlib.util.module_from_spec(_spec)
_spec.loader.exec_module(_mod)
NS, VML_CONTENT = _mod.NS, _mod.VML_CONTENT
zip_read, zip_rewrite, patch_tag_attr, insert_before_close = _mod.zip_read, _mod.zip_rewrite, _mod.patch_tag_attr, _mod.insert_before_close
remove_fragment, assert_wellformed, assert_zip_valid, sha256_of = _mod.remove_fragment, _mod.assert_wellformed, _mod.assert_zip_valid, _mod.sha256_of

ROOT = Path(__file__).resolve().parents[7]
FIXTURE_DIRECTORIES = {
    "set-snapshot": "📸️set-snapshot-applied",
    "set-main-namespace": "🏛️set-main-namespace-applied",
    "set-drawing-namespace": "🎨️set-drawing-namespace-applied",
    "set-relationship-base": "🔗️set-relationship-base-applied",
    "set-conformance-attribute": "🔖️set-conformance-attribute-applied",
    "remove-conformance-attribute": "🏷️remove-conformance-attribute-applied",
    "insert-alternate-content": "🔀️insert-alternate-content-applied",
    "remove-alternate-content": "🚫️remove-alternate-content-applied",
    "insert-vml-part": "🖼️insert-vml-part-applied",
    "remove-vml-part": "🗑️remove-vml-part-applied",
}
PPTX_VERSION = "1.0.2"
ALT_CONTENT = (
    '<mc:AlternateContent xmlns:mc="http://schemas.openxmlformats.org/markup-compatibility/2006">'
    '<mc:Choice Requires="p14"><p:extLst/></mc:Choice><mc:Fallback><p:extLst/></mc:Fallback></mc:AlternateContent>'
)


def base_pptx() -> bytes:
    p = Presentation()
    p.slides.add_slide(p.slide_layouts[6])
    buf = io.BytesIO()
    p.save(buf)
    return buf.getvalue()


def with_presentation_xml(data: bytes, new_xml: str) -> bytes:
    assert_wellformed(new_xml)
    out = zip_rewrite(data, {"ppt/presentation.xml": new_xml})
    assert_zip_valid(out)
    return out


def add_vml_part(data: bytes) -> bytes:
    ct = zip_read(data, "[Content_Types].xml")
    if "vmlDrawing" not in ct:
        ct = ct.replace("</Types>", '<Default Extension="vml" ContentType="application/vnd.openxmlformats-officedocument.vmlDrawing"/></Types>')
    rels = zip_read(data, "ppt/_rels/presentation.xml.rels")
    rels = rels.replace(
        "</Relationships>",
        '<Relationship Id="rIdVml1" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/vmlDrawing" Target="vmlDrawing1.vml"/></Relationships>',
    )
    out = zip_rewrite(data, {"[Content_Types].xml": ct, "ppt/_rels/presentation.xml.rels": rels}, adds={"ppt/vmlDrawing1.vml": VML_CONTENT})
    assert_zip_valid(out)
    return out


def remove_vml_part(data: bytes) -> bytes:
    ct = zip_read(data, "[Content_Types].xml").replace('<Default Extension="vml" ContentType="application/vnd.openxmlformats-officedocument.vmlDrawing"/>', "")
    rels = zip_read(data, "ppt/_rels/presentation.xml.rels").replace(
        '<Relationship Id="rIdVml1" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/vmlDrawing" Target="vmlDrawing1.vml"/>', ""
    )
    out = zip_rewrite(data, {"[Content_Types].xml": ct, "ppt/_rels/presentation.xml.rels": rels}, removes=["ppt/vmlDrawing1.vml"])
    assert_zip_valid(out)
    return out


def build_for_subset(subset: str):
    base = base_pptx()
    tp_uri, strict_p_uri = NS["presentationml"]
    tr_uri, strict_r_uri = NS["relationships"]
    ta_uri, strict_a_uri = NS["drawingml"]
    conf_value = subset

    entries = []

    p2 = Presentation()
    p2.slides.add_slide(p2.slide_layouts[6])
    p2.slides.add_slide(p2.slide_layouts[6])
    buf2 = io.BytesIO()
    p2.save(buf2)
    entries.append(("set-snapshot", "third-party-generated", base, buf2.getvalue(),
                     "Whole-presentation snapshot replace: an unrelated valid presentation (two slides instead of one) substituted wholesale."))

    xml = zip_read(base, "ppt/presentation.xml")

    after_xml = patch_tag_attr(xml, "p:presentation", "conformance", conf_value)
    entries.append(("set-conformance-attribute", "handcrafted", base, with_presentation_xml(base, after_xml),
                     f'The root <p:presentation>\'s conformance attribute set, absent -> "{conf_value}".'))

    before_conf = with_presentation_xml(base, patch_tag_attr(xml, "p:presentation", "conformance", conf_value))
    entries.append(("remove-conformance-attribute", "handcrafted", before_conf, base,
                     f'The root <p:presentation>\'s conformance attribute removed, "{conf_value}" -> absent.'))

    after_ns = patch_tag_attr(xml, "p:presentation", "xmlns:p", strict_p_uri)
    entries.append(("set-main-namespace", "handcrafted", base, with_presentation_xml(base, after_ns),
                     f"The root <p:presentation>'s main xmlns:p namespace URI changed, Transitional ({tp_uri}) -> Strict ({strict_p_uri})."))

    after_r = patch_tag_attr(xml, "p:presentation", "xmlns:r", strict_r_uri)
    entries.append(("set-relationship-base", "handcrafted", base, with_presentation_xml(base, after_r),
                     f"The root <p:presentation>'s relationships xmlns:r namespace URI changed, Transitional ({tr_uri}) -> Strict ({strict_r_uri})."))

    after_a = patch_tag_attr(xml, "p:presentation", "xmlns:a", strict_a_uri)
    entries.append(("set-drawing-namespace", "handcrafted", base, with_presentation_xml(base, after_a),
                     f"The root <p:presentation>'s drawingml xmlns:a namespace URI changed, Transitional ({ta_uri}) -> Strict ({strict_a_uri})."))

    if subset == "strict":
        after_vml = add_vml_part(base)
        entries.append(("insert-vml-part", "handcrafted", base, after_vml,
                         "A legacy VML drawing part (ppt/vmlDrawing1.vml) added to the package, with its Content_Types.xml Default and a presentation.xml.rels relationship wiring it in."))
        entries.append(("remove-vml-part", "handcrafted", after_vml, remove_vml_part(after_vml),
                         "The VML drawing part, its content-type default and its relationship removed, inverse of insert-vml-part."))

        after_alt_xml = insert_before_close(xml, "</p:presentation>", ALT_CONTENT)
        after_alt = with_presentation_xml(base, after_alt_xml)
        entries.append(("insert-alternate-content", "handcrafted", base, after_alt,
                         "An <mc:AlternateContent> markup-compatibility block (Choice/Fallback pair) inserted as the last child of the root <p:presentation>."))
        after_remove_xml = remove_fragment(after_alt_xml, ALT_CONTENT)
        entries.append(("remove-alternate-content", "handcrafted", with_presentation_xml(base, after_alt_xml), with_presentation_xml(base, after_remove_xml),
                         "The <mc:AlternateContent> block removed, inverse of insert-alternate-content."))

    return entries


def emit(subset: str):
    subset_name = {"strict": "🔒️strict", "transitional": "🌉️transitional"}[subset]
    subset_dir = ROOT / f"✏️s/🔌️plugins/🗄️stdio/🗿️artifacts/📽️pptx/🏅️standards/🔖️ecma-376/🪆️subsets/{subset_name}"
    fixtures = subset_dir / "🧫️fixtures"
    oracle_json = subset_dir / "🔮️oracle/🔣️.json"
    reader_oracle = f"python-pptx-pptx-ecma-376-{subset}-mutate-reader"

    manifests = []
    for mutation_id, klass, before_bytes, after_bytes, note in build_for_subset(subset):
        case_dir = fixtures / FIXTURE_DIRECTORIES[mutation_id]
        case_dir.mkdir(parents=True, exist_ok=True)
        (case_dir / "⬅️before.pptx").write_bytes(before_bytes)
        (case_dir / "➡️after.pptx").write_bytes(after_bytes)

        entry = {
            "schema": "semio.repository-test.fixture/v2",
            "id": f"{mutation_id}-applied",
            "class": klass,
            "target": {"artifact": "s.stdio.pptx", "standard": "ecma-376", "subset": subset},
            "mutation": mutation_id,
            "outcome": "applied",
            "units": {"length": "unitless", "angle": "degree"},
            "files": [
                {"role": "expected-before-pptx", "path": f"../🧫️fixtures/{FIXTURE_DIRECTORIES[mutation_id]}/⬅️before.pptx", "mediaType": "application/vnd.openxmlformats-officedocument.presentationml.presentation", "sha256": sha256_of(before_bytes), "bytes": len(before_bytes)},
                {"role": "expected-after-pptx", "path": f"../🧫️fixtures/{FIXTURE_DIRECTORIES[mutation_id]}/➡️after.pptx", "mediaType": "application/vnd.openxmlformats-officedocument.presentationml.presentation", "sha256": sha256_of(after_bytes), "bytes": len(after_bytes)},
            ],
            "provenance": {
                "source": "generated" if klass == "third-party-generated" else "authored",
                "license": "MIT (python-pptx)" if klass == "third-party-generated" else "n/a (handcrafted zip/XML patch of an MIT python-pptx base package)",
                "attribution": "Written by python-pptx 1.0.2's own Presentation.save()" if klass == "third-party-generated" else "A genuine python-pptx 1.0.2 package hand-patched at the OOXML structural level (see notes); the patched ppt/presentation.xml is re-parsed with lxml to confirm well-formedness and the archive's own zip integrity is re-checked before commit",
                "security": "scanned-clean",
                "privacy": "no-personal-data",
            },
            "comparisonProfile": "exact-bytes-v1",
            "reproducible": True,
            "family": "structural",
            "notes": note,
        }
        if klass == "third-party-generated":
            entry["generator"] = {
                "oracle": reader_oracle,
                "packageVersion": PPTX_VERSION,
                "engineFamily": "python-pptx",
                "engineVersion": PPTX_VERSION,
                "command": "uv run python3 🔨️f2-gen-pptx-structural-fixtures.py (python-pptx Presentation object model + .save())",
                "platform": "darwin-arm64",
            }
        manifests.append(entry)
        print(f"[{subset}] {mutation_id:28s} {klass:20s} before={len(before_bytes)}B after={len(after_bytes)}B")

    data = json.loads(oracle_json.read_text())
    data["fixtureManifests"] = manifests
    oracle_json.write_text(json.dumps(data, indent=2, ensure_ascii=False) + "\n")
    print(f"Wrote {len(manifests)} fixtureManifests entries into {oracle_json}\n")


def main() -> None:
    emit("strict")
    emit("transitional")


if __name__ == "__main__":
    main()
