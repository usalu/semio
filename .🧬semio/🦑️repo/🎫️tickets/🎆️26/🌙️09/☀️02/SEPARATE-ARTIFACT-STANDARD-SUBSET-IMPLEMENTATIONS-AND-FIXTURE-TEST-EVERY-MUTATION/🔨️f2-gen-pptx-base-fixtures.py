#!/usr/bin/env python3
"""🔨️ F2 — generates real before/after PPTX fixtures for s.stdio.pptx@ecma-376/base's 7 unfixtured
content-level mutations, all produced through python-pptx 1.0.2's own real object model
(Presentation, slides.add_slide, shapes.add_textbox, text_frame.text, shape.left/.top, the slide-id
list for reordering) -- every byte on disk is python-pptx's own .save() output.
Idempotent: safe to re-run.
"""
import hashlib
import io
import json
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches

ROOT = Path(__file__).resolve().parents[7]
FIXTURE_DIRECTORIES = {
    "insert-slide": "➕️insert-slide-applied",
    "remove-slide": "➖️remove-slide-applied",
    "move-slide": "🔀️move-slide-applied",
    "insert-shape": "🔷️insert-shape-applied",
    "remove-shape": "🔶️remove-shape-applied",
    "set-shape-position": "📐️set-shape-position-applied",
    "set-shape-text": "✍️set-shape-text-applied",
}
SUBSET = ROOT / "✏️s/🔌️plugins/🗄️stdio/🗿️artifacts/📽️pptx/🏅️standards/🔖️ecma-376/🪆️subsets/🧱️base"
FIXTURES = SUBSET / "🧫️fixtures"
ORACLE_JSON = SUBSET / "🔮️oracle/🔣️.json"
READER_ORACLE_ID = "python-pptx-pptx-ecma-376-mutate-reader"
PPTX_VERSION = "1.0.2"


def save(prs: Presentation) -> bytes:
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def one_slide_with_box(text="Hello", left=Inches(1), top=Inches(1)):
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    box = slide.shapes.add_textbox(left, top, Inches(2), Inches(1))
    box.text_frame.text = text
    return prs


def sha256_of(data: bytes) -> str:
    return f"sha256:{hashlib.sha256(data).hexdigest()}"


def main() -> None:
    entries = []

    # 1. insert-slide: 1 slide -> 2 slides.
    prs = Presentation()
    prs.slides.add_slide(prs.slide_layouts[6])
    before = save(prs)
    prs.slides.add_slide(prs.slide_layouts[6])
    after = save(prs)
    entries.append(("insert-slide", before, after, "A second blank slide appended to the presentation."))

    # 2. remove-slide: 2 slides -> 1 slide (inverse).
    prs = Presentation()
    prs.slides.add_slide(prs.slide_layouts[6])
    prs.slides.add_slide(prs.slide_layouts[6])
    before = save(prs)
    xml_slides = prs.slides._sldIdLst
    xml_slides.remove(list(xml_slides)[-1])
    after = save(prs)
    entries.append(("remove-slide", before, after, "The second slide removed from the presentation, inverse of insert-slide."))

    # 3. move-slide: [A, B] -> [B, A], via the slide-id list python-pptx exposes.
    prs = Presentation()
    sa = prs.slides.add_slide(prs.slide_layouts[6])
    sa.shapes.add_textbox(Inches(1), Inches(1), Inches(2), Inches(1)).text_frame.text = "Slide A"
    sb = prs.slides.add_slide(prs.slide_layouts[6])
    sb.shapes.add_textbox(Inches(1), Inches(1), Inches(2), Inches(1)).text_frame.text = "Slide B"
    before = save(prs)
    xml_slides = prs.slides._sldIdLst
    first = list(xml_slides)[0]
    xml_slides.remove(first)
    xml_slides.append(first)
    after = save(prs)
    entries.append(("move-slide", before, after, "Slide order reversed, [A, B] -> [B, A], via the presentation's own slide-id list."))

    # 4. insert-shape: 1 textbox -> 2 textboxes.
    prs = one_slide_with_box("A", Inches(1), Inches(1))
    before = save(prs)
    slide = prs.slides[0]
    slide.shapes.add_textbox(Inches(3), Inches(3), Inches(2), Inches(1)).text_frame.text = "B"
    after = save(prs)
    entries.append(("insert-shape", before, after, "A second textbox shape added to the slide."))

    # 5. remove-shape: 2 textboxes -> 1 textbox (inverse).
    prs = one_slide_with_box("A", Inches(1), Inches(1))
    slide = prs.slides[0]
    second = slide.shapes.add_textbox(Inches(3), Inches(3), Inches(2), Inches(1))
    second.text_frame.text = "B"
    before = save(prs)
    second._element.getparent().remove(second._element)
    after = save(prs)
    entries.append(("remove-shape", before, after, "The second textbox shape removed, inverse of insert-shape."))

    # 6. set-shape-position: shape moved (1in,1in) -> (2in,2in).
    prs = one_slide_with_box("A", Inches(1), Inches(1))
    before = save(prs)
    shape = prs.slides[0].shapes[0]
    shape.left = Inches(2)
    shape.top = Inches(2)
    after = save(prs)
    entries.append(("set-shape-position", before, after, "The shape's left/top offset changed (1in,1in) -> (2in,2in)."))

    # 7. set-shape-text: text content Hello -> World.
    prs = one_slide_with_box("Hello", Inches(1), Inches(1))
    before = save(prs)
    prs.slides[0].shapes[0].text_frame.text = "World"
    after = save(prs)
    entries.append(("set-shape-text", before, after, "The shape's text-frame content replaced, Hello -> World."))

    manifests = []
    for mutation_id, before_bytes, after_bytes, note in entries:
        case_dir = FIXTURES / FIXTURE_DIRECTORIES[mutation_id]
        case_dir.mkdir(parents=True, exist_ok=True)
        (case_dir / "⬅️before.pptx").write_bytes(before_bytes)
        (case_dir / "➡️after.pptx").write_bytes(after_bytes)

        manifests.append({
            "schema": "semio.repository-test.fixture/v2",
            "id": f"{mutation_id}-applied",
            "class": "third-party-generated",
            "target": {"artifact": "s.stdio.pptx", "standard": "ecma-376", "subset": "base"},
            "mutation": mutation_id,
            "outcome": "applied",
            "units": {"length": "unitless", "angle": "degree"},
            "files": [
                {"role": "expected-before-pptx", "path": f"../🧫️fixtures/{FIXTURE_DIRECTORIES[mutation_id]}/⬅️before.pptx", "mediaType": "application/vnd.openxmlformats-officedocument.presentationml.presentation", "sha256": sha256_of(before_bytes), "bytes": len(before_bytes)},
                {"role": "expected-after-pptx", "path": f"../🧫️fixtures/{FIXTURE_DIRECTORIES[mutation_id]}/➡️after.pptx", "mediaType": "application/vnd.openxmlformats-officedocument.presentationml.presentation", "sha256": sha256_of(after_bytes), "bytes": len(after_bytes)},
            ],
            "generator": {
                "oracle": READER_ORACLE_ID,
                "packageVersion": PPTX_VERSION,
                "engineFamily": "python-pptx",
                "engineVersion": PPTX_VERSION,
                "command": "uv run python3 🔨️f2-gen-pptx-base-fixtures.py (python-pptx Presentation object model + .save())",
                "platform": "darwin-arm64",
            },
            "provenance": {
                "source": "generated",
                "license": "MIT (python-pptx)",
                "attribution": "Written by python-pptx 1.0.2's own Presentation.save()",
                "security": "scanned-clean",
                "privacy": "no-personal-data",
            },
            "comparisonProfile": "exact-bytes-v1",
            "reproducible": True,
            "family": "structural",
            "notes": note,
        })
        print(f"{mutation_id:20s} before={len(before_bytes)}B after={len(after_bytes)}B")

    data = json.loads(ORACLE_JSON.read_text())
    data["fixtureManifests"] = manifests
    ORACLE_JSON.write_text(json.dumps(data, indent=2, ensure_ascii=False) + "\n")
    print(f"\nWrote {len(manifests)} fixtureManifests entries into {ORACLE_JSON}")


if __name__ == "__main__":
    main()
